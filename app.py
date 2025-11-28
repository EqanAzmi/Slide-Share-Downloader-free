import os
import re
import io
import json
from datetime import datetime
from urllib.parse import urlparse
from concurrent.futures import ThreadPoolExecutor, as_completed

from flask import Flask, render_template, request, send_file, jsonify, Response
import requests
from bs4 import BeautifulSoup
from PIL import Image
import img2pdf
from pptx import Presentation
from pptx.util import Inches

app = Flask(__name__)
app.secret_key = os.environ.get("SESSION_SECRET", "slideshare-downloader-secret-key")

SITE_DOMAIN = "slidesharedownloaderfree.com"
SITE_URL = f"https://{SITE_DOMAIN}"

SESSION = requests.Session()
SESSION.headers.update({
    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
    'Accept': 'image/webp,image/apng,image/*,*/*;q=0.8',
    'Accept-Language': 'en-US,en;q=0.5',
    'Connection': 'keep-alive',
})

HEADERS = {
    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
    'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
}


def validate_slideshare_url(url):
    if not url:
        return False, "Please provide a URL"
    
    try:
        parsed = urlparse(url)
        valid_domains = ['www.slideshare.net', 'slideshare.net', 'pt.slideshare.net', 
                         'de.slideshare.net', 'es.slideshare.net', 'fr.slideshare.net']
        if parsed.netloc not in valid_domains:
            return False, "Please provide a valid SlideShare URL (https://www.slideshare.net/...)"
        if not parsed.path or parsed.path == '/':
            return False, "Invalid SlideShare presentation URL"
        return True, "Valid URL"
    except Exception as e:
        return False, f"Invalid URL format: {str(e)}"


def extract_slide_images(url):
    try:
        response = requests.get(url, headers=HEADERS, timeout=20)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, 'html.parser')
        
        next_data = soup.select_one('#__NEXT_DATA__')
        if next_data:
            try:
                data = json.loads(next_data.text)
                slideshow = data.get('props', {}).get('pageProps', {}).get('slideshow', {})
                slides = slideshow.get('slides', {})
                total_slides = slideshow.get('totalSlides', 0)
                
                if slides and total_slides > 0:
                    host = slides.get('host', '')
                    image_location = slides.get('imageLocation', '')
                    title = slides.get('title', '')
                    image_sizes = slides.get('imageSizes', [])
                    
                    if host and image_location and title and image_sizes:
                        best_size = image_sizes[-1]
                        quality = best_size.get('quality', 100)
                        width = best_size.get('width', 1280)
                        
                        image_urls = []
                        for i in range(1, total_slides + 1):
                            img_url = f"{host}/{image_location}/{quality}/{title}-{i}-{width}.jpg"
                            image_urls.append(img_url)
                        
                        return image_urls, title, f"Found {len(image_urls)} slides"
            except (json.JSONDecodeError, KeyError):
                pass
        
        image_urls = extract_images_fallback(soup, response.text)
        if image_urls:
            return image_urls, "presentation", f"Found {len(image_urls)} slides"
        
        return None, None, "Could not find slide images. The presentation may be private or SlideShare's format has changed."
        
    except requests.exceptions.Timeout:
        return None, None, "Request timed out. Please try again."
    except requests.exceptions.RequestException as e:
        return None, None, f"Failed to fetch presentation: {str(e)}"
    except Exception as e:
        return None, None, f"Error extracting slides: {str(e)}"


def extract_images_fallback(soup, html_content):
    urls = []
    
    patterns = [
        r'https://image\.slidesharecdn\.com/[^"\'\s]+/\d+/[^"\'\s]+-\d+-\d+\.jpg',
        r'"slideImageUrl"\s*:\s*"([^"]+)"',
    ]
    
    for pattern in patterns:
        matches = re.findall(pattern, html_content)
        for match in matches:
            url = match.replace('\\u002F', '/').replace('\\/', '/')
            if url.startswith('http') and 'slidesharecdn.com' in url:
                urls.append(url)
    
    seen = set()
    cleaned = []
    for url in urls:
        base = re.sub(r'\?.*$', '', url)
        if base not in seen and 'avatar' not in url.lower():
            seen.add(base)
            cleaned.append(url)
    
    def get_slide_num(u):
        match = re.search(r'-(\d+)-\d+\.jpg', u)
        return int(match.group(1)) if match else 0
    
    cleaned.sort(key=get_slide_num)
    return cleaned


def download_single_image(args):
    url, index = args
    try:
        response = SESSION.get(url, timeout=10)
        response.raise_for_status()
        return (index, response.content, None)
    except Exception as e:
        return (index, None, str(e))


def download_images_fast(image_urls):
    results = [None] * len(image_urls)
    
    with ThreadPoolExecutor(max_workers=20) as executor:
        futures = {executor.submit(download_single_image, (url, i)): i 
                   for i, url in enumerate(image_urls)}
        
        for future in as_completed(futures):
            index, img_bytes, error = future.result()
            if img_bytes:
                results[index] = img_bytes
    
    return [b for b in results if b is not None]


def create_pdf_fast(image_bytes_list):
    if not image_bytes_list:
        return None, "No images to convert"
    
    try:
        processed_images = []
        for img_bytes in image_bytes_list:
            img = Image.open(io.BytesIO(img_bytes))
            if img.format == 'WEBP' or img.mode not in ('RGB', 'L'):
                if img.mode in ('RGBA', 'LA', 'P'):
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    if img.mode == 'P':
                        img = img.convert('RGBA')
                    if img.mode in ('RGBA', 'LA'):
                        background.paste(img, mask=img.split()[-1])
                    else:
                        background.paste(img)
                    img = background
                else:
                    img = img.convert('RGB')
                out = io.BytesIO()
                img.save(out, format='JPEG', quality=90)
                processed_images.append(out.getvalue())
            else:
                processed_images.append(img_bytes)
        
        pdf_content = img2pdf.convert(processed_images)
        if pdf_content is None:
            return None, "Failed to convert images to PDF"
        
        pdf_bytes = io.BytesIO()
        pdf_bytes.write(pdf_content)
        pdf_bytes.seek(0)
        
        return pdf_bytes, "PDF created successfully"
        
    except Exception as e:
        return None, f"Failed to create PDF: {str(e)}"


def create_pptx_fast(image_bytes_list):
    if not image_bytes_list:
        return None, "No images to convert"
    
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]
        
        slide_width = 13.333
        slide_height = 7.5
        slide_aspect = slide_width / slide_height
        
        for img_bytes in image_bytes_list:
            slide = prs.slides.add_slide(blank_layout)
            
            img = Image.open(io.BytesIO(img_bytes))
            img_width, img_height = img.size
            img_aspect = img_width / img_height
            
            if img.format == 'WEBP' or img.mode not in ('RGB', 'L'):
                if img.mode in ('RGBA', 'LA', 'P'):
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    if img.mode == 'P':
                        img = img.convert('RGBA')
                    if img.mode in ('RGBA', 'LA'):
                        background.paste(img, mask=img.split()[-1])
                    else:
                        background.paste(img)
                    img = background
                else:
                    img = img.convert('RGB')
            
            img_stream = io.BytesIO()
            img.save(img_stream, format='JPEG', quality=85)
            img_stream.seek(0)
            
            if img_aspect > slide_aspect:
                width = Inches(slide_width)
                height = Inches(slide_width / img_aspect)
                left = Inches(0)
                top = Inches((slide_height - slide_width / img_aspect) / 2)
            else:
                height = Inches(slide_height)
                width = Inches(slide_height * img_aspect)
                left = Inches((slide_width - slide_height * img_aspect) / 2)
                top = Inches(0)
            
            slide.shapes.add_picture(img_stream, left, top, width, height)
        
        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)
        
        return pptx_bytes, "PPTX created successfully"
        
    except Exception as e:
        return None, f"Failed to create PPTX: {str(e)}"


@app.context_processor
def inject_globals():
    return {
        'current_year': datetime.now().year,
        'site_domain': SITE_DOMAIN,
        'site_url': SITE_URL
    }


BLOG_ARTICLES = {
    'how-to-download-slideshare-presentations': {
        'title': 'How to Download SlideShare Presentations: A Complete Guide',
        'excerpt': 'Learn the step-by-step process to download any public SlideShare presentation as PDF or PowerPoint. Includes tips for best quality and ethical use.',
        'date': 'November 27, 2024',
        'read_time': '8 min read',
        'content': '''
            <h2>Introduction to SlideShare Downloads</h2>
            <p>SlideShare has become one of the world's largest platforms for sharing presentations, with millions of educational slideshows, business presentations, and informative content available for viewing. While the platform offers great online viewing capabilities, many users need to download presentations for offline access, study, or reference purposes.</p>
            
            <p>In this comprehensive guide, we'll walk you through everything you need to know about downloading SlideShare presentations safely and effectively using our <a href="/">free SlideShare downloader tool</a>.</p>
            
            <h2>Why Download SlideShare Presentations?</h2>
            <p>There are many legitimate reasons why you might want to download a SlideShare presentation:</p>
            <ul>
                <li><strong>Offline Study:</strong> Access educational content without internet connection</li>
                <li><strong>Research Reference:</strong> Keep presentations for citation and academic research</li>
                <li><strong>Training Materials:</strong> Save helpful training content for future reference</li>
                <li><strong>Backup:</strong> Archive important presentations before they might be removed</li>
                <li><strong>Presentation Prep:</strong> Study presentation styles and techniques</li>
            </ul>
            
            <h2>Step-by-Step Download Guide</h2>
            <h3>Step 1: Find Your Presentation</h3>
            <p>Navigate to <a href="https://www.slideshare.net" target="_blank" rel="noopener">SlideShare.net</a> and find the presentation you want to download. Make sure it's a public presentation - private or login-required presentations cannot be downloaded.</p>
            
            <h3>Step 2: Copy the URL</h3>
            <p>Copy the complete URL from your browser's address bar. It should look something like:</p>
            <code style="display: block; background: #f3f4f6; padding: 1rem; border-radius: 8px; margin: 1rem 0;">https://www.slideshare.net/username/presentation-name</code>
            
            <h3>Step 3: Paste into Our Tool</h3>
            <p>Go to our <a href="/">SlideShare Downloader homepage</a> and paste the URL into the input field.</p>
            
            <h3>Step 4: Choose Your Format</h3>
            <p>Select your preferred download format:</p>
            <ul>
                <li><strong>PDF:</strong> Best for viewing, sharing, and printing. Opens on any device.</li>
                <li><strong>PPTX:</strong> Best if you need to edit the slides in PowerPoint or similar software.</li>
            </ul>
            
            <h3>Step 5: Download</h3>
            <p>Click the download button and wait a few seconds. Your presentation will be automatically downloaded to your device.</p>
            
            <h2>Tips for Best Results</h2>
            <ul>
                <li>Ensure you have a stable internet connection</li>
                <li>Use the complete, unmodified SlideShare URL</li>
                <li>For very large presentations, PDF format may download faster</li>
                <li>If a download fails, try refreshing the page and attempting again</li>
            </ul>
            
            <h2>Ethical Use of Downloaded Content</h2>
            <p>When downloading and using SlideShare presentations, please follow these ethical guidelines:</p>
            <ul>
                <li><strong>Attribution:</strong> Always credit the original author when referencing content</li>
                <li><strong>Personal Use:</strong> Downloaded content should primarily be for personal or educational use</li>
                <li><strong>No Redistribution:</strong> Don't upload downloaded content elsewhere without permission</li>
                <li><strong>Respect Copyright:</strong> Be aware of and respect copyright restrictions</li>
            </ul>
            <p>For more information about copyright and proper use, please review our <a href="/terms">Terms of Service</a> and <a href="/dmca">DMCA Policy</a>.</p>
            
            <h2>Troubleshooting Common Issues</h2>
            <h3>Download Fails</h3>
            <p>If your download fails, it could be due to:</p>
            <ul>
                <li>The presentation is private or requires login</li>
                <li>Temporary server issues - try again in a few minutes</li>
                <li>The presentation has been removed from SlideShare</li>
            </ul>
            
            <h3>Poor Quality Output</h3>
            <p>The quality of your download depends on the original upload quality. Our tool preserves the maximum available quality from SlideShare.</p>
            
            <h2>Conclusion</h2>
            <p>Downloading SlideShare presentations is simple with the right tool. Our <a href="/">free SlideShare downloader</a> makes it easy to save presentations in PDF or PowerPoint format for offline access. Remember to use downloaded content responsibly and respect the original creators' intellectual property rights.</p>
            
            <p>Ready to get started? <a href="/">Try our SlideShare Downloader now</a> - it's fast, free, and requires no registration!</p>
        '''
    }
}


@app.route('/')
def index():
    return render_template('index.html')


@app.route('/blog')
def blog():
    return render_template('blog.html')


@app.route('/blog/<slug>')
def blog_article(slug):
    article = BLOG_ARTICLES.get(slug)
    if not article:
        return render_template('blog.html'), 404
    return render_template('article.html', article=article)


@app.route('/dmca')
def dmca():
    return render_template('dmca.html')


@app.route('/terms')
def terms():
    return render_template('terms.html')


@app.route('/privacy')
def privacy():
    return render_template('privacy.html')


@app.route('/disclaimer')
def disclaimer():
    return render_template('disclaimer.html')


@app.route('/sitemap.xml')
def sitemap():
    try:
        with open('sitemap.xml', 'r') as f:
            content = f.read()
        return Response(content, mimetype='application/xml')
    except FileNotFoundError:
        return "Sitemap not found", 404


@app.route('/robots.txt')
def robots():
    try:
        with open('robots.txt', 'r') as f:
            content = f.read()
        return Response(content, mimetype='text/plain')
    except FileNotFoundError:
        return "Robots.txt not found", 404


@app.route('/download', methods=['POST'])
def download():
    try:
        data = request.get_json()
        url = data.get('url', '').strip()
        format_type = data.get('format', 'pdf').lower()
        
        is_valid, message = validate_slideshare_url(url)
        if not is_valid:
            return jsonify({'success': False, 'error': message}), 400
        
        if format_type not in ['pdf', 'pptx']:
            return jsonify({'success': False, 'error': 'Invalid format. Choose PDF or PPTX'}), 400
        
        image_urls, title, extract_message = extract_slide_images(url)
        if not image_urls:
            return jsonify({'success': False, 'error': extract_message}), 400
        
        image_bytes_list = download_images_fast(image_urls)
        if not image_bytes_list:
            return jsonify({'success': False, 'error': 'Failed to download slide images'}), 400
        
        filename_base = title if title else 'presentation'
        filename_base = re.sub(r'[^\w\-\s]', '_', filename_base.strip())[:100]
        
        if format_type == 'pdf':
            file_bytes, create_message = create_pdf_fast(image_bytes_list)
            if not file_bytes:
                return jsonify({'success': False, 'error': create_message}), 500
            
            return send_file(
                file_bytes,
                mimetype='application/pdf',
                as_attachment=True,
                download_name=f'{filename_base}.pdf'
            )
        else:
            file_bytes, create_message = create_pptx_fast(image_bytes_list)
            if not file_bytes:
                return jsonify({'success': False, 'error': create_message}), 500
            
            return send_file(
                file_bytes,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                as_attachment=True,
                download_name=f'{filename_base}.pptx'
            )
            
    except Exception as e:
        return jsonify({'success': False, 'error': f'An unexpected error occurred: {str(e)}'}), 500


@app.route('/health')
def health():
    return jsonify({'status': 'healthy'})


if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)
